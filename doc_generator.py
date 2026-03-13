"""
doc_generator.py — Motor de generación de documentos para CesarIA
Tipos: docx, pdf, xlsx, pptx + cualquier archivo de código/texto
"""

import io, json, re, os, base64
from datetime import datetime

CODE_EXTENSIONS = {
    'python': 'py', 'py': 'py',
    'javascript': 'js', 'js': 'js',
    'typescript': 'ts', 'ts': 'ts',
    'react': 'jsx', 'jsx': 'jsx', 'tsx': 'tsx',
    'html': 'html', 'css': 'css', 'scss': 'scss',
    'json': 'json', 'sql': 'sql',
    'bash': 'sh', 'shell': 'sh', 'sh': 'sh',
    'yaml': 'yaml', 'yml': 'yaml', 'xml': 'xml',
    'java': 'java', 'cpp': 'cpp', 'c++': 'cpp', 'c': 'c',
    'go': 'go', 'golang': 'go', 'rust': 'rs', 'rs': 'rs',
    'php': 'php', 'ruby': 'rb', 'rb': 'rb', 'swift': 'swift',
    'kotlin': 'kt', 'dart': 'dart',
    'markdown': 'md', 'md': 'md', 'csv': 'csv',
    'txt': 'txt', 'text': 'txt',
    'env': 'env', 'dockerfile': 'dockerfile',
    'toml': 'toml', 'ini': 'ini',
}

MIME_MAP = {
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'pdf':  'application/pdf',
    'py': 'text/x-python', 'js': 'text/javascript', 'ts': 'text/typescript',
    'jsx': 'text/javascript', 'tsx': 'text/typescript',
    'html': 'text/html', 'css': 'text/css', 'json': 'application/json',
    'sql': 'text/x-sql', 'sh': 'text/x-sh', 'yaml': 'text/yaml',
    'xml': 'application/xml', 'java': 'text/x-java',
    'cpp': 'text/x-c++src', 'c': 'text/x-csrc', 'go': 'text/x-go',
    'rs': 'text/x-rust', 'php': 'application/x-httpd-php',
    'rb': 'text/x-ruby', 'swift': 'text/x-swift',
    'kt': 'text/x-kotlin', 'dart': 'text/x-dart',
    'md': 'text/markdown', 'csv': 'text/csv',
    'txt': 'text/plain', 'env': 'text/plain',
    'toml': 'text/plain', 'ini': 'text/plain',
    'dockerfile': 'text/plain',
}


def detect_doc_type(user_message: str) -> dict | None:
    msg = user_message.lower()
    if any(w in msg for w in ['presentación','presentacion','slides','diapositivas','pptx','powerpoint','pitch deck','deck']):
        fname = re.search(r'[\w\-]+\.pptx', msg)
        return {'type': 'pptx', 'filename': fname.group(0) if fname else 'presentacion.pptx'}
    if any(w in msg for w in ['documento word','word','docx','reporte','informe','carta','memo','contrato','propuesta','ensayo','documento']):
        fname = re.search(r'[\w\-]+\.docx', msg)
        return {'type': 'docx', 'filename': fname.group(0) if fname else 'documento.docx'}
    if any(w in msg for w in ['excel','xlsx','hoja de calculo','hoja de cálculo','spreadsheet','presupuesto']):
        fname = re.search(r'[\w\-]+\.xlsx', msg)
        return {'type': 'xlsx', 'filename': fname.group(0) if fname else 'datos.xlsx'}
    if any(w in msg for w in [' pdf',' en pdf','genera pdf','crea pdf']):
        return {'type': 'pdf', 'filename': 'documento.pdf'}
    for keyword, ext in CODE_EXTENSIONS.items():
        patterns = [f'código {keyword}',f'codigo {keyword}',f'script {keyword}',
                    f'archivo {keyword}',f'en {keyword}',f'.{ext}',f'programa {keyword}']
        if any(p in msg for p in patterns):
            fname = re.search(r'[\w\-]+\.' + re.escape(ext), msg)
            return {'type': ext, 'filename': fname.group(0) if fname else f'archivo.{ext}'}
    if any(w in msg for w in ['markdown','.md','readme']):
        return {'type': 'md', 'filename': 'documento.md'}
    return None


def extract_code(response: str, ext: str) -> str:
    lang_map = {
        'py':['python','py'],'js':['javascript','js'],'ts':['typescript','ts'],
        'jsx':['jsx','react'],'tsx':['tsx'],'html':['html','htm'],
        'css':['css','scss'],'json':['json'],'sql':['sql'],
        'sh':['bash','shell','sh'],'yaml':['yaml','yml'],'xml':['xml'],
        'java':['java'],'cpp':['cpp','c++'],'c':['c'],'go':['go','golang'],
        'rs':['rust','rs'],'php':['php'],'rb':['ruby','rb'],'swift':['swift'],
        'md':['markdown','md'],
    }
    aliases = lang_map.get(ext, [ext])
    pattern = r'```(?:' + '|'.join(aliases) + r')?\n([\s\S]*?)```'
    matches = re.findall(pattern, response, re.IGNORECASE)
    if matches:
        return max(matches, key=len).strip()
    all_blocks = re.findall(r'```\w*\n([\s\S]*?)```', response)
    if all_blocks:
        return max(all_blocks, key=len).strip()
    return response.strip()


def generate_file(doc_type: str, content: str, filename: str) -> bytes:
    if doc_type == 'docx':   return _gen_docx(content)
    elif doc_type == 'xlsx': return _gen_xlsx(content)
    elif doc_type == 'pptx': return _gen_pptx(content)
    elif doc_type == 'pdf':  return _gen_pdf(content)
    elif doc_type in ('md','txt','csv','env','toml','ini','dockerfile'):
        return content.encode('utf-8')
    else:
        return extract_code(content, doc_type).encode('utf-8')


# ── DOCX ──────────────────────────────────────────────────────────────────────
def _gen_docx(content: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    doc.styles['Normal'].font.name = 'Calibri'
    doc.styles['Normal'].font.size = Pt(11)
    for s in doc.sections:
        s.top_margin=Cm(2.5); s.bottom_margin=Cm(2.5)
        s.left_margin=Cm(3); s.right_margin=Cm(3)

    def apply_inline(para, text):
        para.clear()
        parts = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                r=para.add_run(part[2:-2]); r.bold=True
            elif part.startswith('*') and part.endswith('*'):
                r=para.add_run(part[1:-1]); r.italic=True
            elif part.startswith('`') and part.endswith('`'):
                r=para.add_run(part[1:-1])
                r.font.name='Courier New'; r.font.size=Pt(10)
                r.font.color.rgb=RGBColor(0x7C,0x6A,0xF7)
            elif part:
                para.add_run(part)

    lines = content.split('\n'); i = 0
    while i < len(lines):
        line = lines[i]
        if line.startswith('### '):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith('## '):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith('# '):
            h=doc.add_heading(line[2:].strip(), level=1)
            h.alignment=WD_ALIGN_PARAGRAPH.CENTER
        elif line.strip() in ('---','==='):
            doc.add_paragraph()
        elif line.startswith('- ') or line.startswith('* '):
            p=doc.add_paragraph(style='List Bullet'); apply_inline(p, line[2:].strip())
        elif re.match(r'^\d+\. ', line):
            p=doc.add_paragraph(style='List Number'); apply_inline(p, re.sub(r'^\d+\. ','',line).strip())
        elif line.startswith('```'):
            code_lines=[]
            i+=1
            while i<len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i]); i+=1
            p=doc.add_paragraph('\n'.join(code_lines))
            p.style=doc.styles['No Spacing']
            for r in p.runs:
                r.font.name='Courier New'; r.font.size=Pt(9)
                r.font.color.rgb=RGBColor(0x2D,0x2D,0x2D)
        elif line.startswith('|') and '|' in line[1:]:
            cells=[c.strip() for c in line.split('|') if c.strip()]
            if not cells or all(set(c)<=set('-: ') for c in cells):
                i+=1; continue
            tbl=doc.add_table(rows=1,cols=len(cells)); tbl.style='Table Grid'
            hdr=tbl.rows[0].cells
            for j,ct in enumerate(cells):
                hdr[j].text=ct
                for r in hdr[j].paragraphs[0].runs:
                    r.bold=True; r.font.color.rgb=RGBColor(255,255,255)
                tc=hdr[j]._tc; tcPr=tc.get_or_add_tcPr()
                shd=OxmlElement('w:shd')
                shd.set(qn('w:fill'),'7C6AF7'); shd.set(qn('w:color'),'auto')
                shd.set(qn('w:val'),'clear'); tcPr.append(shd)
            i+=1
            while i<len(lines) and lines[i].startswith('|'):
                rc=[c.strip() for c in lines[i].split('|') if c.strip()]
                if rc and not all(set(c)<=set('-: ') for c in rc):
                    row=tbl.add_row().cells
                    for j,ct in enumerate(rc[:len(cells)]): row[j].text=ct
                i+=1
            continue
        elif line.strip():
            p=doc.add_paragraph(); apply_inline(p, line.strip())
        else:
            if i>0 and lines[i-1].strip(): doc.add_paragraph()
        i+=1

    buf=io.BytesIO(); doc.save(buf); return buf.getvalue()


# ── XLSX ──────────────────────────────────────────────────────────────────────
def _gen_xlsx(content: str) -> bytes:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb=openpyxl.Workbook(); ws=wb.active; ws.title='Datos'
    ACCENT='FF7C6AF7'; LIGHT='FFF0EFFE'; WHITE='FFFFFFFF'; DARK='FF2D2D4E'
    hdr_font=Font(name='Calibri',bold=True,color=WHITE,size=11)
    hdr_fill=PatternFill('solid',fgColor=ACCENT)
    nor_font=Font(name='Calibri',size=10)
    thin=Border(left=Side(style='thin',color='FFD0C8FF'),right=Side(style='thin',color='FFD0C8FF'),
                top=Side(style='thin',color='FFD0C8FF'),bottom=Side(style='thin',color='FFD0C8FF'))
    c_align=Alignment(horizontal='center',vertical='center')
    l_align=Alignment(horizontal='left',vertical='center',wrap_text=True)
    row=1
    for line in content.split('\n'):
        s=line.strip()
        if not s: row+=1; continue
        if s.startswith('# '):
            c=ws.cell(row=row,column=1,value=s[2:])
            c.font=Font(name='Calibri',bold=True,color=WHITE,size=14)
            c.fill=PatternFill('solid',fgColor=ACCENT)
            c.alignment=c_align
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=8)
            ws.row_dimensions[row].height=28; row+=1
        elif s.startswith('## '):
            c=ws.cell(row=row,column=1,value=s[3:])
            c.font=Font(name='Calibri',bold=True,color=DARK,size=10)
            c.fill=PatternFill('solid',fgColor=LIGHT); c.alignment=l_align
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=8)
            ws.row_dimensions[row].height=20; row+=1
        elif s.startswith('|'):
            cells=[c.strip() for c in s.split('|') if c.strip()]
            if not cells or all(set(c)<=set('-:| ') for c in cells): continue
            prev_val=ws.cell(row=row-1,column=1).value if row>1 else None
            is_hdr=(prev_val is None or str(prev_val or '').startswith('#'))
            for col_idx,cv in enumerate(cells,1):
                c=ws.cell(row=row,column=col_idx,value=cv)
                c.border=thin; c.alignment=l_align
                try:
                    num=float(cv.replace(',','.').replace('%',''))
                    if '%' in cv: c.value=num/100; c.number_format='0.0%'
                    elif '.' in cv: c.value=num; c.number_format='#,##0.00'
                    else: c.value=int(num); c.number_format='#,##0'
                    c.alignment=Alignment(horizontal='right',vertical='center')
                except: pass
                if is_hdr: c.font=hdr_font; c.fill=hdr_fill; c.alignment=c_align
                else:
                    c.font=nor_font
                    if row%2==0: c.fill=PatternFill('solid',fgColor='FFF8F7FF')
            row+=1
        elif s.startswith('- ') or s.startswith('* '):
            c=ws.cell(row=row,column=1,value='• '+s[2:])
            c.font=nor_font; c.alignment=l_align
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=8); row+=1
        elif re.match(r'^\d+\. ',s):
            c=ws.cell(row=row,column=1,value=re.sub(r'^\d+\. ','',s))
            c.font=nor_font; c.alignment=l_align
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=8); row+=1
        else:
            c=ws.cell(row=row,column=1,value=s)
            c.font=nor_font; c.alignment=l_align
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=8); row+=1

    for col in ws.columns:
        ml=max((len(str(c.value)) for c in col if c.value),default=10)
        ws.column_dimensions[get_column_letter(col[0].column)].width=min(max(ml+2,10),50)
    ws.freeze_panes='A2'
    buf=io.BytesIO(); wb.save(buf); return buf.getvalue()


# ── PDF ───────────────────────────────────────────────────────────────────────
def _gen_pdf(content: str) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     Preformatted, Table, TableStyle, HRFlowable)
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

    PURPLE=colors.HexColor('#7C6AF7'); PINK=colors.HexColor('#E040FB')
    DARK=colors.HexColor('#1a1a30'); GRAY=colors.HexColor('#8888AA')
    LGRAY=colors.HexColor('#F0EFFE')
    buf=io.BytesIO()
    doc=SimpleDocTemplate(buf,pagesize=A4,leftMargin=3*cm,rightMargin=3*cm,
                           topMargin=2.5*cm,bottomMargin=2.5*cm)
    styles=getSampleStyleSheet()
    s_h1=ParagraphStyle('H1',parent=styles['Heading1'],fontSize=20,textColor=PURPLE,
                          spaceAfter=12,alignment=TA_CENTER,fontName='Helvetica-Bold')
    s_h2=ParagraphStyle('H2',parent=styles['Heading2'],fontSize=15,textColor=PURPLE,
                          spaceAfter=8,fontName='Helvetica-Bold')
    s_h3=ParagraphStyle('H3',parent=styles['Heading3'],fontSize=12,textColor=DARK,
                          spaceAfter=6,fontName='Helvetica-Bold')
    s_body=ParagraphStyle('Body',parent=styles['Normal'],fontSize=10.5,leading=16,
                           spaceAfter=6,alignment=TA_JUSTIFY,fontName='Helvetica')
    s_bullet=ParagraphStyle('Bullet',parent=s_body,bulletIndent=10,leftIndent=20,spaceAfter=3)
    s_code=ParagraphStyle('Code',fontName='Courier',fontSize=8.5,backColor=LGRAY,
                           leftIndent=10,rightIndent=10,leading=12,spaceAfter=8)

    def md_inline(t):
        t=re.sub(r'\*\*(.+?)\*\*',r'<b>\1</b>',t)
        t=re.sub(r'\*(.+?)\*',r'<i>\1</i>',t)
        t=re.sub(r'`(.+?)`',r'<font name="Courier">\1</font>',t)
        return t

    story=[]; lines=content.split('\n'); i=0
    while i<len(lines):
        line=lines[i]; s=line.strip()
        if s.startswith('### '): story.append(Paragraph(s[4:],s_h3))
        elif s.startswith('## '):
            story.append(HRFlowable(width='100%',thickness=1,color=PURPLE))
            story.append(Paragraph(s[3:],s_h2))
        elif s.startswith('# '):
            story.append(Paragraph(s[2:],s_h1))
            story.append(HRFlowable(width='60%',thickness=2,color=PINK,hAlign='CENTER'))
            story.append(Spacer(1,6))
        elif s.startswith('```'):
            code_lines=[]
            i+=1
            while i<len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i]); i+=1
            story.append(Preformatted('\n'.join(code_lines),s_code))
        elif s.startswith('- ') or s.startswith('* '):
            story.append(Paragraph('• '+md_inline(s[2:]),s_bullet))
        elif re.match(r'^\d+\. ',s):
            story.append(Paragraph(md_inline(re.sub(r'^\d+\. ','',s)),s_bullet))
        elif s.startswith('|'):
            tbl_rows=[]
            while i<len(lines) and lines[i].strip().startswith('|'):
                cells=[c.strip() for c in lines[i].split('|') if c.strip()]
                if cells and not all(set(c)<=set('-:| ') for c in cells):
                    tbl_rows.append(cells)
                i+=1
            if tbl_rows:
                mx=max(len(r) for r in tbl_rows)
                tbl_rows=[r+['']*(mx-len(r)) for r in tbl_rows]
                t=Table(tbl_rows,repeatRows=1)
                t.setStyle(TableStyle([
                    ('BACKGROUND',(0,0),(-1,0),PURPLE),('TEXTCOLOR',(0,0),(-1,0),colors.white),
                    ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('FONTSIZE',(0,0),(-1,-1),9),
                    ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,LGRAY]),
                    ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#D0C8FF')),
                    ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
                    ('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
                ]))
                story.append(t); story.append(Spacer(1,6))
            continue
        elif s=='---':
            story.append(HRFlowable(width='100%',thickness=0.5,color=GRAY))
        elif s: story.append(Paragraph(md_inline(s),s_body))
        else:   story.append(Spacer(1,4))
        i+=1
    doc.build(story)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────────────────────────
def _parse_slides(content: str) -> list:
    sections=re.split(r'\n---+\n|\n#{1,2} Slide \d+[:\.]?',content,flags=re.IGNORECASE)
    if len(sections)<=1:
        parts=re.split(r'\n(?=# )',content)
        sections=parts if len(parts)>1 else [content]
    slides=[]
    for sec in sections:
        sec=sec.strip()
        if not sec: continue
        lines=sec.split('\n'); title=''; subtitle=''; bullets=[]; j=0
        for line in lines:
            s=line.strip()
            if not s: continue
            if not title:
                title=re.sub(r'^#+\s*','',s).strip()
                title=re.sub(r'(?i)^(slide\s*\d+[:\.]?\s*)','',title).strip()
                j+=1; continue
            if not subtitle and not s.startswith('-') and not s.startswith('*') and j==1:
                subtitle=s.strip('*_'); j+=1; continue
            if s.startswith('- ') or s.startswith('* '):
                bullets.append(s[2:].strip())
            elif re.match(r'^\d+\. ',s):
                bullets.append(re.sub(r'^\d+\. ','',s).strip())
            elif s and not s.startswith('#') and len(s)<120:
                bullets.append(s)
            j+=1
        if title: slides.append({'title':title,'subtitle':subtitle,'bullets':bullets})
    if not slides:
        slides=[{'title':'Presentación','subtitle':'',
                 'bullets':[l.strip() for l in content.split('\n') if l.strip() and not l.startswith('#')][:8]}]
    return slides


def _gen_pptx(content: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    C_BG=RGBColor(0x0A,0x0A,0x1A); C_BG2=RGBColor(0x0F,0x0F,0x28)
    C_ACC=RGBColor(0x7C,0x6A,0xF7); C_ACC2=RGBColor(0xE0,0x40,0xFB)
    C_ACC3=RGBColor(0x00,0xE5,0xFF); C_TEXT=RGBColor(0xE8,0xE8,0xF0)
    C_TEXT2=RGBColor(0x88,0x88,0xAA); C_WHT=RGBColor(0xFF,0xFF,0xFF)

    prs=Presentation()
    prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]

    def fill_bg(slide,color):
        bg=slide.background; f=bg.fill; f.solid(); f.fore_color.rgb=color

    def rect(slide,l,t,w,h,color):
        sh=slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background()
        return sh

    def txb(slide,text,l,t,w,h,size=18,bold=False,color=None,align=PP_ALIGN.LEFT,italic=False):
        tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=align
        r=p.add_run(); r.text=text
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color or C_TEXT
        return tb

    def grad_bar(slide,top):
        rect(slide,0,top,6.67,0.06,C_ACC); rect(slide,6.67,top,6.66,0.06,C_ACC2)

    slides_data=_parse_slides(content)
    for idx,sd in enumerate(slides_data):
        slide=prs.slides.add_slide(blank)
        fill_bg(slide,C_BG2)
        title=sd.get('title',''); bullets=sd.get('bullets',[]); subtitle=sd.get('subtitle','')

        if idx==0:
            fill_bg(slide,C_BG)
            rect(slide,0,0,13.33,3.9,RGBColor(0x10,0x10,0x25))
            rect(slide,0,3.84,13.33,0.06,C_ACC)
            circ=slide.shapes.add_shape(9,Inches(10.5),Inches(0.4),Inches(3.0),Inches(3.0))
            circ.fill.solid(); circ.fill.fore_color.rgb=C_ACC; circ.line.fill.background()
            circ2=slide.shapes.add_shape(9,Inches(11.3),Inches(1.2),Inches(1.9),Inches(1.9))
            circ2.fill.solid(); circ2.fill.fore_color.rgb=C_BG; circ2.line.fill.background()
            txb(slide,'⚡',0.4,0.4,1,0.8,size=36,color=C_ACC)
            txb(slide,'CesarIA',1.3,0.5,3,0.6,size=14,bold=True,color=C_ACC)
            txb(slide,title,0.8,1.8,11.2,2.0,size=42,bold=True,color=C_WHT,align=PP_ALIGN.CENTER)
            if subtitle: txb(slide,subtitle,2,4.1,9.33,0.9,size=18,color=C_TEXT2,
                              align=PP_ALIGN.CENTER,italic=True)
            txb(slide,datetime.now().strftime('%B %Y'),0.4,6.8,4,0.4,size=10,color=C_TEXT2)
        else:
            grad_bar(slide,0); grad_bar(slide,7.44)
            txb(slide,str(idx+1),12.4,0.1,0.8,0.4,size=11,color=C_ACC2,align=PP_ALIGN.RIGHT)
            rect(slide,0.5,0.2,12.33,1.05,RGBColor(0x13,0x13,0x2A))
            txb(slide,title,0.7,0.25,11.5,0.9,size=28,bold=True,color=C_WHT)
            rect(slide,0.5,1.45,0.07,5.5,C_ACC)
            if bullets:
                if len(bullets)>4:
                    mid=(len(bullets)+1)//2
                    cols=[bullets[:mid],bullets[mid:]]; cx=[0.75,7.0]; cw=5.8
                else:
                    cols=[bullets]; cx=[0.75]; cw=11.8
                for ci,col_bs in enumerate(cols):
                    y=1.6
                    for bt in col_bs:
                        rect(slide,cx[ci],y,0.3,0.38,RGBColor(0x1E,0x1E,0x35))
                        txb(slide,'◆',cx[ci]+0.01,y+0.01,0.28,0.36,size=11,
                            color=C_ACC,align=PP_ALIGN.CENTER)
                        txb(slide,bt,cx[ci]+0.38,y,cw-0.38,0.45,size=14,color=C_TEXT)
                        y+=0.62
            if subtitle:
                rect(slide,0.5,6.6,12.33,0.7,RGBColor(0x13,0x13,0x2A))
                txb(slide,subtitle,0.7,6.65,12,0.6,size=11,color=C_TEXT2,italic=True)

    buf=io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── PREVIEW DE PPTX (PIL thumbnails) ─────────────────────────────────────────
def pptx_to_preview(pptx_bytes: bytes) -> list:
    """Genera lista de imágenes base64 para preview en browser."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        slides=_parse_slides_from_bytes(pptx_bytes)
        images=[]
        W,H=800,450
        BG=(10,10,26); ACC=(124,106,247); ACC2=(224,64,251)
        WHT=(232,232,240); GRY=(136,136,170); BG2=(15,15,40)

        FONT_PATHS=[
            '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
            '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
            '/usr/share/fonts/truetype/freefont/FreeSansBold.ttf',
        ]
        FONT_PATHS_NB=[
            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
            '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
        ]
        def load_font(paths,size):
            for p in paths:
                try: return ImageFont.truetype(p,size)
                except: pass
            return ImageFont.load_default()

        f_title=load_font(FONT_PATHS,26); f_body=load_font(FONT_PATHS_NB,15)
        f_small=load_font(FONT_PATHS_NB,11)

        for idx,sd in enumerate(slides):
            img=Image.new('RGB',(W,H),BG)
            draw=ImageDraw.Draw(img)
            # Barras degradadas
            for x in range(W):
                t=x/W
                r=int(ACC[0]*(1-t)+ACC2[0]*t); g=int(ACC[1]*(1-t)+ACC2[1]*t); b=int(ACC[2]*(1-t)+ACC2[2]*t)
                draw.line([(x,0),(x,5)],fill=(r,g,b))
                draw.line([(x,H-5),(x,H)],fill=(r,g,b))
            draw.rectangle([30,18,W-30,92],fill=(19,19,42))
            draw.rectangle([30,18,39,92],fill=ACC)
            title=sd.get('title','')
            draw.text((50,26),title[:52],fill=WHT,font=f_title)
            if idx>0:
                draw.text((W-45,22),str(idx+1),fill=ACC2,font=f_small)
            y=105
            for bt in sd.get('bullets',[])[:7]:
                draw.rectangle([50,y+5,57,y+12],fill=ACC)
                draw.text((68,y),bt[:72],fill=WHT,font=f_body)
                y+=40
            if sd.get('subtitle'):
                draw.rectangle([30,H-58,W-30,H-14],fill=(19,19,42))
                draw.text((48,H-50),sd['subtitle'][:90],fill=GRY,font=f_small)
            draw.text((30,H-16),'⚡ CesarIA',fill=ACC,font=f_small)
            b64buf=io.BytesIO(); img.save(b64buf,format='JPEG',quality=88)
            images.append('data:image/jpeg;base64,'+base64.b64encode(b64buf.getvalue()).decode())
        return images
    except Exception as e:
        print(f'[PREVIEW] ❌ {e}')
        return []


def _parse_slides_from_bytes(pptx_bytes: bytes) -> list:
    from pptx import Presentation
    prs=Presentation(io.BytesIO(pptx_bytes)); slides=[]
    for slide in prs.slides:
        title=''; bullets=[]
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text=shape.text_frame.text.strip()
            if not text: continue
            if not title: title=text
            elif text!=title:
                for line in text.split('\n'):
                    l=line.strip()
                    if l and l!=title: bullets.append(l)
        slides.append({'title':title or 'Slide','bullets':bullets[:8],'subtitle':''})
    return slides
